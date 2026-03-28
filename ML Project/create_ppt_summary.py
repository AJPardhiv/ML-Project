"""
create_ppt_summary.py

Builds a 1-slide PPT-ready summary using generated results and plots.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
VIS_DIR = BASE_DIR / "visualizations"
METRICS_PATH = VIS_DIR / "metrics_summary.json"
OUTPUT_PPTX = BASE_DIR / "Gesture_Project_One_Page_Summary.pptx"


def load_metrics() -> dict:
    if not METRICS_PATH.exists():
        raise FileNotFoundError(
            f"Missing metrics file: {METRICS_PATH}. Run generate_analysis_report.py first."
        )
    return json.loads(METRICS_PATH.read_text(encoding="utf-8"))


def add_title(slide, text: str):
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.55))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(22)
    p.alignment = PP_ALIGN.CENTER


def add_summary_block(slide, metrics: dict):
    box = slide.shapes.add_textbox(Inches(0.3), Inches(0.75), Inches(6.2), Inches(1.45))
    tf = box.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "Dataset"
    p0.runs[0].font.bold = True
    p0.runs[0].font.size = Pt(14)

    p1 = tf.add_paragraph()
    p1.text = "• 4600 samples, 63 landmark features, 4 gesture classes"
    p1.level = 1
    p1.runs[0].font.size = Pt(11)

    p2 = tf.add_paragraph()
    p2.text = "• Classes: move, click, scroll, pause (near-balanced)"
    p2.level = 1
    p2.runs[0].font.size = Pt(11)

    p3 = tf.add_paragraph()
    p3.text = "Model"
    p3.runs[0].font.bold = True
    p3.runs[0].font.size = Pt(14)

    p4 = tf.add_paragraph()
    p4.text = (
        "• RandomForest (n_estimators=100, max_depth=15, "
        "min_samples_split=5, min_samples_leaf=2)"
    )
    p4.level = 1
    p4.runs[0].font.size = Pt(11)

    metrics_box = slide.shapes.add_textbox(Inches(6.75), Inches(0.75), Inches(6.0), Inches(1.45))
    mtf = metrics_box.text_frame
    mtf.word_wrap = True

    m0 = mtf.paragraphs[0]
    m0.text = "Key Results"
    m0.runs[0].font.bold = True
    m0.runs[0].font.size = Pt(14)

    metric_lines = [
        f"• Test Accuracy: {metrics['test_accuracy']:.4f}",
        f"• Weighted Precision: {metrics['test_precision_weighted']:.4f}",
        f"• Weighted Recall: {metrics['test_recall_weighted']:.4f}",
        f"• Weighted F1-score: {metrics['test_f1_weighted']:.4f}",
        f"• ROC-AUC (micro): {metrics['roc_micro_auc']:.4f}",
        f"• PR-AUC (micro): {metrics['pr_micro_auc']:.4f}",
    ]

    for line in metric_lines:
        p = mtf.add_paragraph()
        p.text = line
        p.level = 1
        p.runs[0].font.size = Pt(11)


def add_caption(slide, left: float, top: float, width: float, text: str):
    cap = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.25))
    ctf = cap.text_frame
    ctf.clear()
    p = ctf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(9)


def maybe_add_image(slide, path: Path, left: float, top: float, width: float, height: float):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def add_plots(slide):
    plot_specs = [
        (VIS_DIR / "oob_error_curve.png", 0.3, 2.35, 3.1, 2.2, "OOB Error Curve"),
        (VIS_DIR / "roc_curve_multiclass.png", 3.45, 2.35, 3.1, 2.2, "ROC Curves"),
        (VIS_DIR / "precision_recall_curve_multiclass.png", 6.6, 2.35, 3.1, 2.2, "Precision-Recall Curves"),
        (VIS_DIR / "confusion_matrix.png", 9.75, 2.35, 3.1, 2.2, "Confusion Matrix"),
    ]

    for img_path, left, top, width, height, caption in plot_specs:
        maybe_add_image(slide, img_path, left, top, width, height)
        add_caption(slide, left, top + 2.22, width, caption)


def add_footer(slide):
    footer = slide.shapes.add_textbox(Inches(0.35), Inches(4.78), Inches(12.5), Inches(0.35))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Generated from sample dataset (gesture_data.csv) | Pipeline: MediaPipe landmarks -> RandomForest classifier"
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(9)


def main() -> None:
    metrics = load_metrics()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Gesture-Based Laptop Controller: One-Page Model Summary")
    add_summary_block(slide, metrics)
    add_plots(slide)
    add_footer(slide)

    prs.save(OUTPUT_PPTX)
    print(f"Saved PPT summary: {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
